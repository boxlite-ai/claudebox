#!/usr/bin/env python3
"""
RisingWave PPT Template Generator

Generates a professional PowerPoint template with RisingWave brand identity.
Brand colors derived from RisingWave's official branding:
- Primary: Deep navy blue (#1E3A8A) - trustworthy, mature
- Secondary: Medium blue (#3B82F6) - accent elements
- Accent: Vibrant green (#0ABD59) - from official GitHub presence
- Wave-inspired decorative elements reflecting the ocean tides logo motif
"""

import math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


# ── Brand Colors ──────────────────────────────────────────────────────────────
DEEP_NAVY = RGBColor(0x1E, 0x3A, 0x8A)       # #1E3A8A - primary brand color
MEDIUM_BLUE = RGBColor(0x3B, 0x82, 0xF6)      # #3B82F6 - secondary blue
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)       # #DBEAFE - light background
ACCENT_GREEN = RGBColor(0x0A, 0xBD, 0x59)     # #0ABD59 - vibrant green accent
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)         # #1F2937 - dark text
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)       # #F3F4F6 - light gray bg
MID_GRAY = RGBColor(0x6B, 0x72, 0x80)         # #6B7280 - secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)             # #FFFFFF
FLUORESCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)  # #8B5CF6 - community color
FLUORESCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)    # #06B6D4 - community color

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_wave_decoration(slide, y_offset=0, color=None, opacity_pct=15):
    """Add wave-shaped decorative elements at the bottom of a slide."""
    if color is None:
        color = MEDIUM_BLUE

    # Create a series of overlapping ellipses to simulate a wave pattern
    wave_y = Inches(6.0) + Emu(y_offset)
    num_waves = 8
    wave_width = Inches(3.5)
    wave_height = Inches(1.8)
    step = Inches(1.5)

    for i in range(num_waves):
        left = Emu(int(step * i) - int(Inches(0.5)))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left, wave_y, wave_width, wave_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.fill.fore_color.brightness = 0.7 if i % 2 == 0 else 0.8
        shape.line.fill.background()
        # Move to back
        sp = shape._element
        sp.getparent().insert(2, sp)


def add_top_bar(slide, color=DEEP_NAVY, height=Inches(0.08)):
    """Add a thin accent bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bottom_bar(slide, color=ACCENT_GREEN, height=Inches(0.06)):
    """Add a thin accent bar at the bottom of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), SLIDE_HEIGHT - height, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_side_stripe(slide, color=DEEP_NAVY, width=Inches(0.5)):
    """Add a vertical stripe on the left side of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), width, SLIDE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 font_color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=DARK_TEXT, bullet_color=MEDIUM_BLUE, font_name="Calibri"):
    """Add a bulleted list text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_before = Pt(8)
        p.space_after = Pt(4)
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.findall(qn('a:buNone'))
        for bn in buNone:
            pPr.remove(bn)
        # Set bullet character
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'})
        buClr.append(srgbClr)
        buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '120000'})
        pPr.append(buClr)
        pPr.append(buSzPct)
        pPr.append(buChar)

        # Add indent
        from lxml import etree
        pPr.set('marL', str(Emu(Inches(0.4))))
        pPr.set('indent', str(Emu(Inches(-0.25))))

    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=14, font_color=WHITE, corner_size=50000):
    """Add a rounded rectangle shape with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Set corner rounding
    shape.adjustments[0] = corner_size / 100000

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_circle(slide, left, top, size, fill_color, text="", font_size=14,
               font_color=WHITE):
    """Add a circle shape with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left, top, size, size
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_page_number(slide, num, total=None):
    """Add page number to bottom right."""
    text = f"{num}" if total is None else f"{num} / {total}"
    add_text_box(
        slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.35),
        text, font_size=10, font_color=MID_GRAY,
        alignment=PP_ALIGN.RIGHT
    )


# ═══════════════════════════════════════════════════════════════════════════════
# Slide Builders
# ═══════════════════════════════════════════════════════════════════════════════

def build_title_slide(prs):
    """Slide 1: Title / Cover Page."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide, DEEP_NAVY)

    # Wave decoration at bottom
    add_wave_decoration(slide, y_offset=Emu(Inches(0.5)), color=MEDIUM_BLUE)

    # Accent line at top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Company name / logo placeholder area
    add_text_box(
        slide, Inches(1.5), Inches(1.0), Inches(4), Inches(0.6),
        "RISINGWAVE", font_size=20, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Decorative accent dot
    add_circle(slide, Inches(1.5), Inches(1.85), Inches(0.12), ACCENT_GREEN)

    # Title
    add_text_box(
        slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.5),
        "Presentation Title Goes Here",
        font_size=44, font_color=WHITE, bold=True, font_name="Calibri"
    )

    # Subtitle
    add_text_box(
        slide, Inches(1.5), Inches(3.9), Inches(8), Inches(0.8),
        "Event Streaming Platform for Agents, Apps, and Analytics",
        font_size=20, font_color=LIGHT_BLUE, font_name="Calibri"
    )

    # Divider line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(5.0), Inches(2.5), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Presenter info
    add_text_box(
        slide, Inches(1.5), Inches(5.3), Inches(5), Inches(0.4),
        "Presenter Name  |  Title  |  Date",
        font_size=14, font_color=LIGHT_BLUE, font_name="Calibri"
    )


def build_agenda_slide(prs):
    """Slide 2: Agenda / Table of Contents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "AGENDA", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(6), Inches(0.7),
        "Today's Topics",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Agenda items with numbered circles
    agenda_items = [
        ("Overview & Introduction", "Brief overview of today's presentation"),
        ("Key Insights & Data", "Important metrics and findings"),
        ("Technical Deep Dive", "Architecture and implementation details"),
        ("Roadmap & Next Steps", "Future plans and action items"),
    ]

    colors = [DEEP_NAVY, MEDIUM_BLUE, FLUORESCENT_CYAN, ACCENT_GREEN]

    for i, (title, desc) in enumerate(agenda_items):
        y = Inches(2.0 + i * 1.2)

        # Number circle
        add_circle(slide, Inches(1.0), y, Inches(0.55), colors[i],
                   text=str(i + 1), font_size=18, font_color=WHITE)

        # Title
        add_text_box(
            slide, Inches(1.8), y - Inches(0.05), Inches(5), Inches(0.4),
            title, font_size=20, font_color=DEEP_NAVY, bold=True
        )

        # Description
        add_text_box(
            slide, Inches(1.8), y + Inches(0.35), Inches(5), Inches(0.4),
            desc, font_size=14, font_color=MID_GRAY
        )

    # Right side decorative element
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    sp = shape._element
    sp.getparent().insert(2, sp)

    # Decorative wave circles on right
    for i, c in enumerate([DEEP_NAVY, MEDIUM_BLUE, FLUORESCENT_CYAN]):
        add_circle(
            slide, Inches(9.5 + i * 0.8), Inches(3.2),
            Inches(1.5), c
        )

    add_page_number(slide, 2)


def build_content_slide(prs):
    """Slide 3: Content with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "SECTION TITLE", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(8), Inches(0.7),
        "Content Slide Title",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Subtitle
    add_text_box(
        slide, Inches(1.0), Inches(1.7), Inches(8), Inches(0.5),
        "Supporting description text that provides context for the content below.",
        font_size=16, font_color=MID_GRAY, font_name="Calibri"
    )

    # Bullet list
    items = [
        "First key point with supporting details and context",
        "Second important insight or finding to highlight",
        "Third actionable item or recommendation for the audience",
        "Fourth point demonstrating depth and thoroughness",
        "Fifth conclusion or call to action for next steps",
    ]
    add_bullet_list(
        slide, Inches(1.0), Inches(2.5), Inches(6.5), Inches(4.0),
        items, font_size=18, font_color=DARK_TEXT, bullet_color=MEDIUM_BLUE
    )

    # Right side accent box
    add_rounded_rect(
        slide, Inches(8.5), Inches(2.5), Inches(4.0), Inches(3.5),
        LIGHT_BLUE
    )

    # Highlight text inside accent box
    add_text_box(
        slide, Inches(8.8), Inches(2.8), Inches(3.4), Inches(0.4),
        "KEY HIGHLIGHT", font_size=12, font_color=MEDIUM_BLUE,
        bold=True, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(8.8), Inches(3.3), Inches(3.4), Inches(0.8),
        "42%", font_size=48, font_color=DEEP_NAVY,
        bold=True, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(8.8), Inches(4.3), Inches(3.4), Inches(1.0),
        "Performance improvement with real-time streaming processing",
        font_size=14, font_color=MID_GRAY, font_name="Calibri"
    )

    add_page_number(slide, 3)


def build_two_column_slide(prs):
    """Slide 4: Two-column layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "COMPARISON", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.7),
        "Two Column Layout",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Left column header
    add_rounded_rect(
        slide, Inches(1.0), Inches(2.0), Inches(5.5), Inches(0.6),
        DEEP_NAVY, text="Before", font_size=16, font_color=WHITE
    )

    # Left column content
    left_items = [
        "Complex batch processing pipelines",
        "High latency data processing",
        "Manual scaling and management",
        "Fragmented tooling ecosystem",
    ]
    add_bullet_list(
        slide, Inches(1.0), Inches(2.8), Inches(5.5), Inches(3.5),
        left_items, font_size=16, font_color=DARK_TEXT, bullet_color=DEEP_NAVY
    )

    # Right column header
    add_rounded_rect(
        slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.6),
        ACCENT_GREEN, text="After", font_size=16, font_color=WHITE
    )

    # Right column content
    right_items = [
        "Simple SQL-based stream processing",
        "Sub-100ms real-time processing",
        "Auto-scaling cloud-native architecture",
        "Unified streaming + lakehouse platform",
    ]
    add_bullet_list(
        slide, Inches(7.0), Inches(2.8), Inches(5.5), Inches(3.5),
        right_items, font_size=16, font_color=DARK_TEXT, bullet_color=ACCENT_GREEN
    )

    # Divider line between columns
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55), Inches(2.0), Inches(0.03), Inches(4.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GRAY
    shape.line.fill.background()

    add_page_number(slide, 4)


def build_metrics_slide(prs):
    """Slide 5: Key metrics / KPIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "KEY METRICS", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.7),
        "Performance at a Glance",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Metrics cards
    metrics = [
        ("<100ms", "Processing Latency", "Real-time stream processing", DEEP_NAVY),
        ("10x", "Cost Reduction", "Compared to traditional pipelines", MEDIUM_BLUE),
        ("99.99%", "Uptime SLA", "Enterprise-grade reliability", FLUORESCENT_CYAN),
        ("1000+", "Customers", "Across global industries", ACCENT_GREEN),
    ]

    for i, (value, label, desc, color) in enumerate(metrics):
        x = Inches(1.0 + i * 3.0)
        y = Inches(2.2)

        # Card background
        card = add_rounded_rect(
            slide, x, y, Inches(2.6), Inches(3.8), LIGHT_GRAY
        )

        # Color accent bar at top of card
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.0), y, Inches(2.6), Inches(0.08)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Metric value
        add_text_box(
            slide, x + Inches(0.2), y + Inches(0.5), Inches(2.2), Inches(1.0),
            value, font_size=40, font_color=color, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Metric label
        add_text_box(
            slide, x + Inches(0.2), y + Inches(1.6), Inches(2.2), Inches(0.5),
            label, font_size=16, font_color=DEEP_NAVY, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Metric description
        add_text_box(
            slide, x + Inches(0.2), y + Inches(2.2), Inches(2.2), Inches(1.0),
            desc, font_size=12, font_color=MID_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    add_page_number(slide, 5)


def build_architecture_slide(prs):
    """Slide 6: Architecture / Diagram placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "ARCHITECTURE", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.7),
        "System Architecture",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Architecture flow boxes
    flow_items = [
        ("Data Sources", "Kafka, Pulsar,\nPostgreSQL CDC,\nIoT Devices", DEEP_NAVY),
        ("Ingestion", "Continuous\nStream Ingestion\n& Validation", MEDIUM_BLUE),
        ("Processing", "SQL-based\nStream Processing\n& Transformation", FLUORESCENT_CYAN),
        ("Materialization", "Materialized Views\n& Real-time\nServing", ACCENT_GREEN),
    ]

    for i, (title, desc, color) in enumerate(flow_items):
        x = Inches(0.8 + i * 3.1)
        y = Inches(2.2)

        # Box
        box = add_rounded_rect(
            slide, x, y, Inches(2.7), Inches(3.2), color
        )

        # Title
        add_text_box(
            slide, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.5),
            title, font_size=18, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Separator
        sep = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.4), y + Inches(1.0), Inches(1.9), Inches(0.02)
        )
        sep.fill.solid()
        sep.fill.fore_color.rgb = WHITE
        sep.line.fill.background()

        # Description
        add_text_box(
            slide, x + Inches(0.2), y + Inches(1.2), Inches(2.3), Inches(1.5),
            desc, font_size=14, font_color=WHITE,
            alignment=PP_ALIGN.CENTER
        )

        # Arrow between boxes (except last)
        if i < len(flow_items) - 1:
            arrow_x = x + Inches(2.7)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x, Inches(3.5), Inches(0.4), Inches(0.35)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

    # Bottom label
    add_text_box(
        slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.5),
        "RisingWave: Unified Streaming + Lakehouse Platform with Apache Iceberg Integration",
        font_size=14, font_color=MID_GRAY, alignment=PP_ALIGN.CENTER
    )

    add_page_number(slide, 6)


def build_quote_slide(prs):
    """Slide 7: Quote / Testimonial slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)

    # Accent bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Large quote mark
    add_text_box(
        slide, Inches(1.5), Inches(1.0), Inches(2), Inches(1.5),
        "\u201C", font_size=120, font_color=MEDIUM_BLUE,
        bold=True, font_name="Georgia"
    )

    # Quote text
    add_text_box(
        slide, Inches(2.0), Inches(2.5), Inches(9), Inches(2.0),
        "RisingWave fundamentally changed how we think about\n"
        "real-time data processing. What used to take days\n"
        "now happens in milliseconds.",
        font_size=28, font_color=WHITE, font_name="Calibri",
        line_spacing=1.5
    )

    # Divider
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(2.0), Inches(5.0), Inches(1.5), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Attribution
    add_text_box(
        slide, Inches(2.0), Inches(5.3), Inches(5), Inches(0.4),
        "Speaker Name", font_size=18, font_color=WHITE,
        bold=True, font_name="Calibri"
    )
    add_text_box(
        slide, Inches(2.0), Inches(5.7), Inches(5), Inches(0.4),
        "Title, Company", font_size=14, font_color=LIGHT_BLUE,
        font_name="Calibri"
    )

    # Wave decoration
    add_wave_decoration(slide, y_offset=Emu(Inches(0.5)), color=MEDIUM_BLUE)

    add_page_number(slide, 7)


def build_team_slide(prs):
    """Slide 8: Team / About slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "OUR TEAM", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.7),
        "Meet the Team",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Team member placeholders
    members = [
        ("Team Member 1", "Role / Title", DEEP_NAVY),
        ("Team Member 2", "Role / Title", MEDIUM_BLUE),
        ("Team Member 3", "Role / Title", FLUORESCENT_CYAN),
        ("Team Member 4", "Role / Title", ACCENT_GREEN),
    ]

    for i, (name, role, color) in enumerate(members):
        x = Inches(1.0 + i * 3.0)
        y = Inches(2.2)

        # Avatar circle placeholder
        circle = add_circle(
            slide, x + Inches(0.55), y, Inches(1.5), LIGHT_GRAY,
            text=name[0], font_size=36, font_color=color
        )

        # Color accent dot
        add_circle(slide, x + Inches(1.6), y + Inches(1.1), Inches(0.3), color)

        # Name
        add_text_box(
            slide, x, y + Inches(1.8), Inches(2.6), Inches(0.4),
            name, font_size=16, font_color=DEEP_NAVY, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Role
        add_text_box(
            slide, x, y + Inches(2.2), Inches(2.6), Inches(0.4),
            role, font_size=13, font_color=MID_GRAY,
            alignment=PP_ALIGN.CENTER
        )

        # Short bio placeholder
        add_text_box(
            slide, x, y + Inches(2.7), Inches(2.6), Inches(1.0),
            "Brief description of experience and expertise area.",
            font_size=11, font_color=MID_GRAY,
            alignment=PP_ALIGN.CENTER
        )

    add_page_number(slide, 8)


def build_chart_slide(prs):
    """Slide 9: Chart / Data visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "DATA & INSIGHTS", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(8), Inches(0.7),
        "Growth & Performance",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Add a sample bar chart
    chart_data = CategoryChartData()
    chart_data.categories = ['Q1', 'Q2', 'Q3', 'Q4']
    chart_data.add_series('Revenue', (120, 180, 250, 340))
    chart_data.add_series('Users', (80, 150, 220, 310))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0), Inches(2.0), Inches(7.0), Inches(4.5),
        chart_data
    )

    # Style the chart
    chart_obj = chart.chart
    chart_obj.has_legend = True
    chart_obj.legend.include_in_layout = False

    # Color the series
    plot = chart_obj.plots[0]
    series_colors = [DEEP_NAVY, ACCENT_GREEN]
    for i, series in enumerate(plot.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = series_colors[i]

    # Right side insights
    add_text_box(
        slide, Inches(8.5), Inches(2.2), Inches(4), Inches(0.5),
        "KEY TAKEAWAYS", font_size=14, font_color=MEDIUM_BLUE,
        bold=True, font_name="Calibri"
    )

    insights = [
        "Revenue grew 183% YoY",
        "User base expanded 3.9x",
        "Consistent quarter-over-quarter growth",
        "Strong product-market fit indicators",
    ]
    add_bullet_list(
        slide, Inches(8.5), Inches(2.8), Inches(4.0), Inches(3.5),
        insights, font_size=14, font_color=DARK_TEXT, bullet_color=ACCENT_GREEN
    )

    add_page_number(slide, 9)


def build_timeline_slide(prs):
    """Slide 10: Timeline / Roadmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide)
    add_bottom_bar(slide)

    # Section label
    add_text_box(
        slide, Inches(1.0), Inches(0.5), Inches(3), Inches(0.4),
        "ROADMAP", font_size=12, font_color=ACCENT_GREEN,
        bold=True, font_name="Calibri"
    )

    # Title
    add_text_box(
        slide, Inches(1.0), Inches(0.9), Inches(10), Inches(0.7),
        "Product Roadmap",
        font_size=36, font_color=DEEP_NAVY, bold=True, font_name="Calibri"
    )

    # Timeline horizontal line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(3.7), Inches(10.5), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MEDIUM_BLUE
    line.line.fill.background()

    # Timeline points
    milestones = [
        ("Q1 2026", "Phase 1", "Foundation &\nCore Features", DEEP_NAVY),
        ("Q2 2026", "Phase 2", "Advanced\nIntegrations", MEDIUM_BLUE),
        ("Q3 2026", "Phase 3", "Scale &\nOptimization", FLUORESCENT_CYAN),
        ("Q4 2026", "Phase 4", "Enterprise\nFeatures", ACCENT_GREEN),
    ]

    for i, (quarter, phase, desc, color) in enumerate(milestones):
        x = Inches(1.8 + i * 2.8)

        # Circle on timeline
        add_circle(slide, x + Inches(0.3), Inches(3.45), Inches(0.55), color,
                   text=str(i + 1), font_size=16, font_color=WHITE)

        # Quarter label (above)
        add_text_box(
            slide, x - Inches(0.1), Inches(2.4), Inches(1.5), Inches(0.4),
            quarter, font_size=14, font_color=MID_GRAY, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Phase label
        add_text_box(
            slide, x - Inches(0.1), Inches(2.8), Inches(1.5), Inches(0.4),
            phase, font_size=12, font_color=color, bold=True,
            alignment=PP_ALIGN.CENTER
        )

        # Description (below)
        add_text_box(
            slide, x - Inches(0.3), Inches(4.2), Inches(2.0), Inches(1.0),
            desc, font_size=14, font_color=DARK_TEXT,
            alignment=PP_ALIGN.CENTER
        )

    add_page_number(slide, 10)


def build_closing_slide(prs):
    """Slide 11: Thank You / Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DEEP_NAVY)

    # Wave decoration
    add_wave_decoration(slide, y_offset=Emu(Inches(0.5)), color=MEDIUM_BLUE)

    # Accent bar at top
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Company name
    add_text_box(
        slide, Inches(0), Inches(1.5), SLIDE_WIDTH, Inches(0.5),
        "RISINGWAVE", font_size=18, font_color=ACCENT_GREEN,
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # Thank you
    add_text_box(
        slide, Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(1.2),
        "Thank You", font_size=56, font_color=WHITE,
        bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri"
    )

    # Divider
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.7), Inches(3.6), Inches(2.0), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()

    # Contact info
    add_text_box(
        slide, Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.5),
        "Questions & Discussion", font_size=24, font_color=LIGHT_BLUE,
        alignment=PP_ALIGN.CENTER
    )

    # Contact details
    add_text_box(
        slide, Inches(0), Inches(5.0), SLIDE_WIDTH, Inches(0.4),
        "risingwave.com  |  github.com/risingwavelabs  |  contact@risingwave.com",
        font_size=14, font_color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER
    )


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def generate_template(output_path="risingwave_template.pptx"):
    """Generate the complete RisingWave PPT template."""
    prs = Presentation()

    # Set widescreen 16:9 dimensions
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Build all slides
    print("Building RisingWave PPT template...")

    print("  [1/11] Title slide...")
    build_title_slide(prs)

    print("  [2/11] Agenda slide...")
    build_agenda_slide(prs)

    print("  [3/11] Content slide...")
    build_content_slide(prs)

    print("  [4/11] Two-column layout...")
    build_two_column_slide(prs)

    print("  [5/11] Metrics slide...")
    build_metrics_slide(prs)

    print("  [6/11] Architecture slide...")
    build_architecture_slide(prs)

    print("  [7/11] Quote slide...")
    build_quote_slide(prs)

    print("  [8/11] Team slide...")
    build_team_slide(prs)

    print("  [9/11] Chart slide...")
    build_chart_slide(prs)

    print("  [10/11] Timeline / Roadmap slide...")
    build_timeline_slide(prs)

    print("  [11/11] Closing slide...")
    build_closing_slide(prs)

    # Save
    prs.save(output_path)
    print(f"\nTemplate saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("\nSlide overview:")
    print("  1. Title / Cover")
    print("  2. Agenda / Table of Contents")
    print("  3. Content with Bullet Points")
    print("  4. Two-Column Comparison")
    print("  5. Key Metrics / KPIs")
    print("  6. Architecture / Diagram")
    print("  7. Quote / Testimonial")
    print("  8. Team / About")
    print("  9. Chart / Data Visualization")
    print("  10. Timeline / Roadmap")
    print("  11. Thank You / Closing")


if __name__ == "__main__":
    generate_template()
